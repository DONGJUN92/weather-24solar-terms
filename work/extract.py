import io, sys, json
from pptx import Presentation
from pptx.util import Emu
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

p = Presentation(sys.argv[1] if len(sys.argv) > 1 else 'work/source.pptx')
IN = 914400.0


def col(r):
    try:
        c = r.font.color
        if c is None or c.type is None:
            return None
        if c.type == 1:  # scheme
            return 'theme:%s' % c.theme_color
        return str(c.rgb)
    except Exception:
        return None


for i, s in enumerate(p.slides, 1):
    print('\n' + '=' * 100)
    print('SLIDE %02d   layout=%s' % (i, s.slide_layout.name))
    print('=' * 100)
    for sh in s.shapes:
        kind = sh.shape_type
        try:
            L, T, W, H = sh.left / IN, sh.top / IN, sh.width / IN, sh.height / IN
            geo = '(%.2f,%.2f %.2fx%.2f in)' % (L, T, W, H)
        except Exception:
            geo = '(no geom)'
        head = '  [%s] name=%r %s' % (kind, sh.name, geo)
        if not sh.has_text_frame:
            print(head + '  <no text>')
            continue
        tf = sh.text_frame
        txt = tf.text.strip()
        print(head)
        if not txt:
            print('       (empty)')
            continue
        for pi, para in enumerate(tf.paragraphs):
            if not para.text.strip():
                continue
            runs = []
            for r in para.runs:
                runs.append({
                    't': r.text,
                    'f': r.font.name,
                    'sz': (r.font.size.pt if r.font.size else None),
                    'b': r.font.bold, 'i': r.font.italic,
                    'c': col(r),
                })
            print('       p%d lvl=%s align=%s' % (pi, para.level, para.alignment))
            for r in runs:
                print('          %-58r f=%-22s sz=%-6s b=%-5s c=%s'
                      % (r['t'][:58], r['f'], r['sz'], r['b'], r['c']))
