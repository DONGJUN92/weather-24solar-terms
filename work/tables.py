import io, sys
from pptx import Presentation
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
p = Presentation('work/source.pptx')
for i, s in enumerate(p.slides, 1):
    for sh in s.shapes:
        if not sh.has_table:
            continue
        t = sh.table
        print('\n=== SLIDE %02d  table %r  %d rows x %d cols ===' % (i, sh.name, len(t.rows), len(t.columns)))
        for ri, row in enumerate(t.rows):
            cells = []
            for c in row.cells:
                txt = c.text.replace('\n', ' / ').strip()
                cells.append(txt[:44])
            print('  r%-2d | %s' % (ri, ' | '.join(cells)))
        # 첫 셀 서식
        c0 = t.cell(0, 0)
        for para in c0.text_frame.paragraphs:
            for r in para.runs:
                print('     [r0c0 서식] font=%s sz=%s bold=%s' % (r.font.name, r.font.size.pt if r.font.size else None, r.font.bold))
                break
            break
print('\n=== SLIDE 13 shapes ===')
s13 = p.slides[12]
for sh in s13.shapes:
    txt = sh.text_frame.text.strip()[:120] if sh.has_text_frame else '<no tf>'
    print('  %-18s %r  ->  %r' % (sh.shape_type, sh.name, txt))
print('\n=== SLIDE 01 shapes (layout placeholders) ===')
s1 = p.slides[0]
for sh in s1.shapes:
    print('  %-18s %r idx=%s type=%s' % (sh.shape_type, sh.name,
          getattr(sh.placeholder_format, 'idx', None) if sh.is_placeholder else None,
          getattr(sh.placeholder_format, 'type', None) if sh.is_placeholder else None))
print('\n--- layout of S01 ---')
for sh in s1.slide_layout.placeholders:
    print('  layout ph idx=%s type=%s  text=%r' % (sh.placeholder_format.idx, sh.placeholder_format.type, sh.text_frame.text[:70]))
