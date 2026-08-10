# -*- coding: utf-8 -*-
"""주최측 템플릿에 내용을 채워 본선 발표자료를 만든다.

입력  assets/deck/템플릿_원본.pptx · assets/deck/*.png (배포본에서 추출한 앱 차트)
출력  발표자료_Weather24_신동준.pptx  (PDF 는 PowerPoint 로 내보낸다)

유지하는 것 (심사위원이 배점표와 연결해 듣도록)
  · 상단 헤더 장식 · 눈썹(평가 항목) · 슬라이드 제목 · ●● 소제목 · 쪽번호 · 마스터/레이아웃

바꾼 것 (템플릿이 '디자인 자유 변경 가능'이라 명시)
  · 본문 박스를 '짧은 개조식 + 그 주장을 증명하는 그림' 구조로 재구성
  · 상세 근거는 발표자 노트로 내렸다 — 슬라이드는 읽는 문서가 아니다
  · 그림은 슬라이드가 주장하는 바로 그것만 넣는다(앱 차트는 실물, 개념·인용은 네이티브 도형)
"""
import copy
import io
import os
import sys

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck_graphics import (AMBER, CREAM, FILL, GREY, NAVY, TERRA, WHITE,  # noqa: E402
                           caption, ladder, matrix2x2, mini_table, quote_card,
                           stat_row, step_strip, textbox)

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(BASE, 'assets', 'deck')
URL = 'https://weather-24solar-terms.vercel.app'
DECK = '발표자료_Weather24_신동준'   # 제출 파일명 — 발표자명 포함


# ───────────────────────────────────────────── 템플릿 서식 보존 편집
def no_bullet(tf):
    for p in tf.paragraphs:
        pPr = p._p.get_or_add_pPr()
        for tag in ('a:buChar', 'a:buAutoNum', 'a:buNone'):
            for el in pPr.findall(qn(tag)):
                pPr.remove(el)
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))


def set_lines(tf, lines, size=None):
    """첫 런을 제자리 편집해 폰트·색을 물려받는다 (text 대입은 서식을 날린다)."""
    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.add_run()
    p0.runs[0].text = lines[0]
    if size:
        p0.runs[0].font.size = Pt(size)
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    for _ in lines[1:]:
        p0._p.getparent().append(copy.deepcopy(p0._p))
    for p, line in zip(list(tf.paragraphs)[1:], lines[1:]):
        p.runs[0].text = line
        if size:
            p.runs[0].font.size = Pt(size)
    if size:
        no_bullet(tf)


def by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError('%r 없음: %s' % (name, [s.name for s in slide.shapes]))


def geo(shape):
    IN = 914400.0
    return (shape.left / IN, shape.top / IN, shape.width / IN, shape.height / IN)


def drop_slide(prs, idx):
    lst = prs.slides._sldIdLst
    items = list(lst)
    rId = items[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    lst.remove(items[idx])


# 발표 대본은 scripts/deck_script.py 한 곳에만 둔다.
# 여기서 노트를 직접 쓰면, build_deck.py 를 다시 돌리는 순간 대본이 짧은 옛 노트로
# 덮여 사라진다 — 실제로 그 상태였다. 저장 직후 apply_notes.apply() 로 주입한다.


def pic(slide, png, cx, y, w=None, h=None):
    """앱 실화면 PNG. cx = 가로 중심(도형 박스 중앙 정렬용). 흰 배경이 채움색 위에
    떠 보이지 않게 카드처럼 얇은 테두리를 준다."""
    from PIL import Image
    from pptx.dml.color import RGBColor
    path = os.path.join(SHOTS, png if png.endswith('.png') else png + '.png')
    if not os.path.exists(path):
        print('  !! 이미지 없음: %s' % png)
        return None
    iw, ih = Image.open(path).size
    if w and not h:
        h = w * ih / iw
    elif h and not w:
        w = h * iw / ih
    p = slide.shapes.add_picture(path, Inches(cx - w / 2.0), Inches(y),
                                 Inches(w), Inches(h))
    p.line.color.rgb = RGBColor(0xC8, 0xD4, 0xE4)
    p.line.width = Pt(0.75)
    return p


# ═════════════════════════════════════════════════════════ 본문
def main():
    prs = Presentation(os.path.join(SHOTS, '템플릿_원본.pptx'))
    S = prs.slides

    # ── 표지 ───────────────────────────────────────────────
    cover = {10: 'Weather24 — 절기, 아직 맞을까',
             11: '팀 Weather24 (1인 팀)',
             12: URL,
             13: '신동준'}
    for sh in S[0].shapes:
        if sh.is_placeholder and sh.placeholder_format.idx in cover:
            set_lines(sh.text_frame, [cover[sh.placeholder_format.idx]])
    print('S1 표지')

    # ── I. 결과물 및 팀 기본 정보 ──────────────────────────
    s = S[4]
    tbl = by_name(s, '표 2').table
    for r, v in enumerate(['팀 Weather24 (1인 팀)',
                           'Weather24 — 절기, 아직 맞을까',
                           URL,
                           '중·고생이 ‘덥다’ 기준을 정해 24절기를 실측 검증하는 기후 학습 웹앱']):
        set_lines(tbl.cell(r, 1).text_frame, [v])
    set_lines(by_name(s, 'TextBox 6').text_frame,
              ['신동준 / 기획·개발·데이터·디자인', '※ 1인 팀 — 단독 수행'])
    # 우측 하단: 이 앱이 무엇을 시키는지 한 그림
    textbox(s, 7.85, 3.62, 4.30, 0.24, ['학습 흐름 — 한 미션의 뼈대'], size=10.5,
            color=TERRA, bold=True)
    step_strip(s, 7.85, 3.92, 4.30, 0.92,
               [('①', '예측 봉인', '자료 비공개'),
                ('②', '직접 조작', '기준·지역·절기'),
                ('③', '내 결론', '주장·근거·추론·한계')])
    stat_row(s, 7.85, 5.05, 4.30, 0.62,
             [('16', '관측지점'), ('24', '절기'), ('3', '지표'), ('5', '미션')])
    caption(s, 7.85, 5.76, 4.30, '기상청 ASOS 1969–2026 실측 · 조작 즉시 재계산')
    print('S5 기본정보')

    # ── 1. 학습 대상 분석 ──────────────────────────────────
    s = S[5]
    set_lines(by_name(s, '텍스트 개체 틀 9').text_frame, ['1. 학습 대상 분석'])
    b = by_name(s, '직사각형 8')
    bx, by_, bw, bh = geo(b)
    set_lines(b.text_frame, [
        '■ 학습 대상 — 중3~고2',
        '  · 교육과정상 ‘기후변화와 지구환경’·‘자료 해석’ 성취기준 배치 학년',
        '  · 임계값 개념과 평균·빈도 구분이 가능한 최소 학년',
        '  · 24절기를 들어는 봤으나 천문 기준임은 미학습 — 오개념 교정 여지 최대',
        '',
        '■ 난이도 적합성 — 같은 개념을 3단 용어로 반복 (아래)',
    ], size=12)
    ladder(s, bx + 0.30, by_ + 1.42, bw - 0.60, 1.42,
           [('생활어', ['‘덥다’ · ‘여름’ — 정의 없이 시작',
                     '등장 위치 · 첫 화면 · 예측 봉인']),
            ('조작어', ['기준 25°C · 기준 이상 더위일 · 마지막 초과일',
                     '등장 위치 · 기준선 조작 · 판정 화면']),
            ('학술어', ['임계값 · 유효 열용량 매개변수 · 계절 지연',
                     '등장 위치 · 열관성 실험실 · 내 결론'])])
    caption(s, bx + 0.30, by_ + 2.90, bw - 0.60,
            '한 개념을 세 번 다른 말로 만나는 구조 · 수식은 끝까지 비노출')
    textbox(s, bx + 0.30, by_ + 3.26, bw - 0.60, 0.24,
            ['■ 필수 경로는 짧게, 심화는 선택 분기로'], size=12, color=NAVY)
    step_strip(s, bx + 0.30, by_ + 3.58, bw - 0.60, 0.86,
               [('필수', '5개 미션', '한 화면 한 질문'),
                ('선택', '열관성 실험실', '에너지수지 모형'),
                ('선택', '전국 16지점', '26개 비교 구간'),
                ('교사', '학습지', '오개념 표 · 루브릭')])
    caption(s, bx + 0.30, by_ + 4.52, bw - 0.60,
            '필수 경로만으로 한 차시 완결 · 심화는 남는 시간과 학생 수준에 맞춰 선택')
    print('S6 학습대상')

    # ── 2. 문제 정의 및 분석 ───────────────────────────────
    s = S[6]
    set_lines(by_name(s, '텍스트 개체 틀 9').text_frame, ['2. 문제 정의 및 분석'])
    L = by_name(s, '직사각형 7'); lx, ly, lw, lh = geo(L)
    set_lines(L.text_frame, [
        '■ 문제 — 절기·날씨·기후의 혼용',
        '  · “처서가 더워졌다” = 천문 날짜에 기온을 귀속하는 오개념',
        '  · 기존 콘텐츠는 결론만 전달 — 학습자의 검증 경험 없음',
        '',
        '■ 개발계획서 제시 개념 → 결과물 반영',
        '  · 24절기의 과학(태양황경 15°) → 궤도 화면 (아래)',
        '  · 기후평년과 계절 길이 → 미션 2 ‘여름은 며칠일까’',
        '  · 기온 상승과 계절 이동 → 미션 1 처서 · 미션 5 계절 지연',
        '  · 습도·강수 변화 → 지표 3종 · 미션 4 빈도 대 강도',
        '  · (확장) CO₂–기온 상관 → ‘지구 맥락’ 화면',
    ], size=12)
    pic(s, 'orbit', lx + 2.16, ly + 2.06, w=3.20)
    textbox(s, lx + 3.94, ly + 2.52, 1.58, 1.00,
            ['처서 (미션 1의 검증 대상)', '= 태양황경 150°',
             '= 양력 8/23 무렵 고정', '↔ 그 무렵 기온은 이동'],
            size=9.5, color=TERRA, space=2)
    caption(s, lx + 0.30, ly + 4.58, lw - 0.60,
            '태양황경 15° 간격으로 정한 천문 날짜 — 기온 기준 아님')
    R = by_name(s, '직사각형 10'); rx, ry, rw, rh = geo(R)
    set_lines(R.text_frame, [
        '■ 주제를 체험 중에 인식하는 세 지점',
        '  · 첫 화면 표제가 곧 주제',
        '  · 모든 조작 화면에 자료 범위 고지',
        '  · 개념 렌즈 4문항 통과 후에야 관측 진입',
    ], size=12)
    quote_card(s, rx + 0.26, ry + 1.18, rw - 0.52, 0.66, '첫 화면 표제',
               ['“24절기, 지금도 맞을까?”',
                '움직이지 않는 절기와 움직이는 기후를 나란히'], accent=AMBER)
    quote_card(s, rx + 0.26, ry + 1.98, rw - 0.52, 1.06, '모든 조작 화면 하단 고지',
               ['“◈ 기상청 ASOS 실측 · 과거 5년(1969–1973) vs',
                '   현재 5년(2021–2025) — 관측 신호이고 30년',
                '   기후평년이 아닙니다 · 절기는 태양 위치로 정한',
                '   천문 날짜라 해마다 거의 움직이지 않습니다”'])
    quote_card(s, rx + 0.26, ry + 3.18, rw - 0.52, 1.32, '개념 렌즈 — 절기·날씨·기후 분류 4문항',
               ['· “처서는 양력 8월 23일 무렵이다” → 절기',
                '· “처서가 지났는데도 어제 낮 기온 31°C” → 날씨',
                '· “서울 한 해 평균기온이 …로 높아졌다” → 기후',
                '· “처서를 9월로 옮겨야 한다” → 절기 (오개념 판별)'], accent=TERRA)
    caption(s, rx + 0.26, ry + 4.56, rw - 0.52, '화면에 실제로 표시되는 문구를 그대로 인용')
    print('S7 문제정의')

    # ── 3. 아이디어 도출 ───────────────────────────────────
    s = S[7]
    set_lines(by_name(s, '텍스트 개체 틀 9').text_frame, ['3. 아이디어 도출'])
    L = by_name(s, '직사각형 12'); lx, ly, lw, lh = geo(L)
    set_lines(L.text_frame, [
        '■ 차별화 — 실측과 모형을 한 앱에서',
        '  · 24절기(문화 기준선)를 과학 검증 대상으로 전환',
        '  · ‘덥다’의 정의를 학습자에게 이양',
        '  · 관측으로 본 것을 물리 모형으로 다시 확인',
    ], size=12)
    matrix2x2(s, lx + 0.24, ly + 1.26, lw - 0.48, 2.66,
              '실측 관측 자료 보유 →', '변수 조작 학습 →',
              {(0, 0): ('PhET · En-ROADS', '모형만 · 실측 대조 없음'),
               (1, 0): ('Weather24', '실측 + 모형 + 문화 기준선'),
               (0, 1): ('일반 기후 교재', '결론만 전달'),
               (1, 1): ('기후 스트라이프 · 개방포털', '실측만 · 조작 학습 없음')})
    caption(s, lx + 0.24, ly + 4.02, lw - 0.48,
            '두 축을 함께 가진 칸의 공백 — 그 자리를 채운 설계')
    textbox(s, lx + 0.24, ly + 4.38, lw - 0.48, 0.48,
            ['■ 문화 기준선을 쓰는 이유 — 학습자가 이미 알고 있어',
             '   “틀렸을 수도 있다”는 반증 동기가 즉시 생김'], size=12, color=NAVY)
    R = by_name(s, '직사각형 13'); rx, ry, rw, rh = geo(R)
    set_lines(R.text_frame, [
        '■ 새로운 학습 경험 — 예측 봉인 루프',
        '  · 자료 공개 전 내 생각 봉인',
        '  · 조작 후 봉인한 예측 재대조',
    ], size=12)
    ys = ry + 1.02
    for i, (num, title, body, acc) in enumerate([
        ('①', '예측 봉인', '자료·그래프 비공개 상태로 내 생각 선택', GREY),
        ('②', '직접 조작', '기준선·지역·절기·지표를 바꿔 확인', NAVY),
        ('③', '봉인 해제', '“내 예측 vs 자료” 대조 — 어긋난 지점 명시', TERRA),
        ('④', '내 결론', '주장·근거·추론·한계를 내 말로 작성', NAVY),
        ('⑤', '이해 확인 → 전이', '오답은 자료로 되돌림 · 다른 절기에 적용', NAVY),
    ]):
        quote_card(s, rx + 0.26, ys + i * 0.66, rw - 0.52, 0.58,
                   num + ' ' + title, [body], accent=acc,
                   label_size=10.5, body_size=10)
    caption(s, rx + 0.26, ys + 5 * 0.66 + 0.04, rw - 0.52,
            '③이 이 앱의 핵심 — 봉인에서 끝내지 않는 필수 대조')
    print('S8 아이디어')

    # ── 4. 체험·참여형 설계 ────────────────────────────────
    s = S[8]
    set_lines(by_name(s, '텍스트 개체 틀 9').text_frame, ['4. 체험∙참여형 설계'])
    L = by_name(s, '직사각형 11'); lx, ly, lw, lh = geo(L)
    set_lines(L.text_frame, [
        '■ 기준선 이동 → 아래 숫자 즉시 재계산 (배포본 실측)',
        '  · 조작 경로 4종 — 손잡이 끌기 / 슬라이더 / ＋− / 자주 쓰는 기준',
        '  · 기준선 미조작 시 판정 버튼 잠금 — 클릭만으로 통과 불가',
    ], size=12)
    mini_table(s, lx + 0.22, ly + 1.00, lw - 0.44,
               [['내가 정한 ‘덥다’', '더위일 과거 → 현재', '더위가 그치는 날'],
                ['22°C', '75일 → 110일', '9/17 → 9/29'],
                ['25°C', '31일 → 68일', '8/30 → 9/12'],
                ['28°C', '4.8일 → 33일', '8/8 → 8/29'],
                ['30°C', '0.6일 → 7.8일', '7/19 → 8/10'],
                ['!33°C', '!0일 → 0일', '!비교 불가 → 안내']],
               colw=[1.5, 1.5, 1.5], row_h=0.27, head_h=0.30, size=10)
    caption(s, lx + 0.22, ly + 3.00, lw - 0.44,
            '서울·기온 실측 — 같은 자료, 기준에 따라 달라지는 결론')
    textbox(s, lx + 0.22, ly + 3.36, lw - 0.44, 1.20, [
        '■ 조작 결과 → 학습 개념 연결',
        '  · 기준↑ → 일수↓ → “모호한 말은 기준을 정해야 자료가 된다”',
        '  · 1mm 감소·50mm 증가 → “빈도와 강도는 다른 지표”',
        '  · 온실효과↑ → 곡선 상승·지연 불변 → “둘은 다른 원인”',
    ], size=12, color=NAVY)
    R = by_name(s, '직사각형 16'); rx, ry, rw, rh = geo(R)
    set_lines(R.text_frame, [
        '■ 탐구 범위 — 16지점 × 24절기 × 3지표 × 26구간',
        '  · 보기 3종 — 그래프 / 전국 지도 / 표',
        '  · 조작부마다 화면 라벨 명시 · 결과 수치를 조작부 바로 위에',
    ], size=12)
    pic(s, 'map16', rx + rw / 2.0, ry + 0.86, h=3.58)
    caption(s, rx + 0.24, ry + 4.52, rw - 0.48,
            '‘지도 16지점’ 보기 — 한 지역의 결과를 전국으로 넓힐 수 없음의 시각적 확인')
    print('S9 체험설계')

    # ── 5. 학습 성과와 학습적 피드백 ───────────────────────
    s = S[9]
    set_lines(by_name(s, '텍스트 개체 틀 9').text_frame, ['5. 학습 성과와 학습적 피드백'])
    b = by_name(s, '직사각형 12'); bx, by_, bw, bh = geo(b)
    set_lines(b.text_frame, [
        '■ 학습 성과 — 완료 화면 ‘내 기록’으로 교사가 회수 (채점 근거)',
    ], size=12)
    stat_row(s, bx + 0.26, by_ + 0.52, bw - 0.52, 0.66,
             [('예측 5건', '봉인 후 대조'), ('CERL 5편', '주장·근거·추론·한계'),
              ('렌즈 4문항', '절기·날씨·기후'), ('이해확인 5문항', '1차/2차 구분'),
              ('전이 3문항', '미학습 상황 적용'), ('실험실 값', '내가 고른 물리량')])
    textbox(s, bx + 0.26, by_ + 1.34, bw - 0.52, 0.26,
            ['■ 잘못된 조작·선택에 대한 학습적 피드백 6종 — 전부 배포본 실측'],
            size=12, color=NAVY)
    cards = [
        ('① 기준을 극단으로 올렸을 때', ['“과거·현재 모두 0일에 가까워 비교할 것이',
                                 '  남지 않아요. 기준선을 내리면 비교가 시작됩니다”'], TERRA),
        ('② 이해 확인 1차 오답', ['“아직이에요. 정답을 알려 드리기 전에',
                             '  한 번 더 볼까요?” + 되돌림 힌트 — 정답 비공개'], TERRA),
        ('③ 오답 되돌림 힌트', ['“기준선을 25°C와 28°C에 놓았을 때 같은 자료인데',
                            '  여름 길이가 달라졌던 것을 떠올려 보세요”'], NAVY),
        ('④ 근거에 단위 누락', ['“그 숫자가 무엇을 센 값인지 단위나 기간을',
                            '  함께 써 주세요 (예: …일, …°C, …년)”'], NAVY),
        ('⑤ 한계에 범위어 누락', ['“어디까지 통하는지가 들어가야 해요 —',
                             '  지역·기간·기준·원인 중 무엇을 넘어 말할 수 없는지”'], NAVY),
        ('⑥ 베끼기·복사 차단', ['“화면의 안내 문장을 그대로 옮긴 것 같아요”',
                           '“‘근거’ 칸과 같은 문장이에요”'], GREY),
    ]
    cw = (bw - 0.52 - 0.24 * 2) / 3.0
    for i, (label, body, acc) in enumerate(cards):
        cx = bx + 0.26 + (i % 3) * (cw + 0.24)
        cy = by_ + 1.68 + (i // 3) * 1.06
        quote_card(s, cx, cy, cw, 0.94, label, body, accent=acc,
                   label_size=10.5, body_size=9.5)
    caption(s, bx + 0.26, by_ + 3.86, bw - 0.52,
            '공통 원칙 — 오답 시 정답 비공개 · 자료로 되돌린 뒤 재질문')
    textbox(s, bx + 0.26, by_ + 4.18, bw - 0.52, 0.46, [
        '■ 우회 차단 — 문장 중복·발문 베끼기 검사 · 오개념(범위 확대·인과 단정) 점검을 필수 경로에서 실행',
    ], size=12, color=NAVY)
    print('S10 학습성과')

    # ── 6. 구현 완성도 ─────────────────────────────────────
    s = S[10]
    set_lines(by_name(s, '텍스트 개체 틀 9').text_frame, ['6. 구현 완성도'])
    L = by_name(s, '직사각형 7'); lx, ly, lw, lh = geo(L)
    set_lines(L.text_frame, [
        '■ 기상청 ASOS 종관기상관측 일자료 1969–2026',
        '  · 16지점 × 24절기 × 3지표(기온·강수·습도)',
        '  · 파생 집계 — 기준 초과 일수 / 마지막 초과일 / 절기 무렵 15일 평균',
        '  · 조작 시 즉시 재계산 — 사전 계산 이미지 아님 (아래 실화면)',
    ], size=12)
    pic(s, 'chart-temp', lx + lw / 2.0, ly + 1.44, w=lw - 0.56)
    caption(s, lx + 0.28, ly + 3.62, lw - 0.56,
            '서울 · 과거(회색 점선) 대 현재(빨강) · 보라색이 내가 정한 기준선')
    stat_row(s, lx + 0.28, ly + 3.94, lw - 0.56, 0.62,
             [('8,360건', '자동 회귀 검증'), ('0건', '겹침·잘림·오류'),
              ('134KB', '첫 로드'), ('0개', '외부 도메인')])
    R = by_name(s, '직사각형 10'); rx, ry, rw, rh = geo(R)
    set_lines(R.text_frame, [
        '■ 제품 안의 AI — 증거 감사관의 과장 차단',
    ], size=12)
    quote_card(s, rx + 0.26, ry + 0.62, rw - 0.52, 0.62, '학생이 쓴 결론 (일부러 과장)',
               ['“서울에서 더위가 늦어졌으니 우리나라 전체가',
                '  기후변화 때문에 완전히 바뀌었다”'], accent=GREY,
               label_size=10.5, body_size=10)
    quote_card(s, rx + 0.26, ry + 1.40, rw - 0.52, 1.02, 'AI 응답 (배포본 실측 · 200 OK)',
               ['“서울 자료만으로 우리나라 전체를 말할 수 없고,',
                '  ‘기후변화 때문에’라는 원인도 이 근거만으로',
                '  단정할 수 없다”',
                '→ 다음 행동: 한계 문장 쓰기'], accent=TERRA,
               label_size=10.5, body_size=10)
    textbox(s, rx + 0.26, ry + 2.60, rw - 0.52, 0.72, [
        '  · 구조화 출력 강제 · 외부 전송 동의 후에만 요청',
        '  · 실패 시 기기 안 규칙 점검으로 폴백 — 학습은 항상 완결',
    ], size=12, color=NAVY)
    textbox(s, rx + 0.26, ry + 3.32, rw - 0.52, 0.26,
            ['■ 제작 과정의 AI — Claude Code'], size=12, color=NAVY)
    stat_row(s, rx + 0.26, ry + 3.62, rw - 0.52, 0.62,
             [('6회', 'AI 레드팀 자체 감사'), ('100/100', 'AI 품질 평가 세트'),
              ('6건', '프롬프트 세션 원문')])
    caption(s, rx + 0.26, ry + 4.32, rw - 0.52,
            '자기 결과물을 AI로 반증 → 전량 수정 → 게이트로 재발 방지')
    print('S11 구현완성도')


    # ── 가이드 3장 + 별지1 삭제 ────────────────────────────
    for idx in sorted([11, 3, 2, 1], reverse=True):
        drop_slide(prs, idx)
    print('가이드 3장 + 별지1 삭제 → 총 %d장' % len(prs.slides._sldIdLst))

    out = os.path.join(BASE, DECK + '.pptx')
    prs.save(out)
    print('\n저장: %s' % out)

    # 발표 대본 주입 — 대본은 scripts/deck_script.py 한 곳에만 있고,
    # apply_notes.py 와 같은 함수를 써서 두 경로가 절대 어긋나지 않게 한다.
    from apply_notes import apply
    print('\n발표자 노트 주입')
    apply(out)


if __name__ == '__main__':
    main()
