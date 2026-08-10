# -*- coding: utf-8 -*-
"""발표자 노트에 발표 대본을 넣는다.

대본은 scripts/deck_script.py 한 곳에만 있고, 두 경로가 모두 그것을 쓴다.
  · scripts/build_deck.py   — 자료를 처음부터 다시 만들 때
  · scripts/apply_notes.py  — 이미 있는 pptx 에 대본만 갈아 끼울 때 (기본 동작)

이렇게 나눈 이유: 사용자가 PowerPoint 에서 저장한 파일(쪽번호 캐시가 고쳐진 판)을
버리지 않고 대본만 교체할 수 있어야 한다.

사용법
  python scripts/apply_notes.py                     # 기본 파일에 적용
  python scripts/apply_notes.py 다른파일.pptx        # 지정 파일에 적용
"""
import io
import os
import sys

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck_script import SCRIPT  # noqa: E402

if __name__ == '__main__':          # import 로 들어올 때 부모의 stdout 을 닫아 버리면 안 된다
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT = os.path.join(BASE, '발표자료_Weather24_신동준.pptx')

# 한국어 발표 속도 — 또박또박 읽을 때의 실측 기준. 시간 추정에만 쓴다.
CHARS_PER_MIN = 330


def set_notes(slide, text):
    """노트 텍스트프레임을 통째로 교체한다. 문단은 줄 단위로 나눈다."""
    tf = slide.notes_slide.notes_text_frame
    lines = text.split('\n')
    tf.text = lines[0]
    for ln in lines[1:]:
        tf.add_paragraph().text = ln


def apply(path):
    prs = Presentation(path)
    n = len(prs.slides)
    if n != len(SCRIPT):
        raise SystemExit('슬라이드 %d장인데 대본은 %d장분입니다' % (n, len(SCRIPT)))
    total = 0
    for i, (slide, block) in enumerate(zip(prs.slides, SCRIPT), 1):
        set_notes(slide, block)
        chars = len(block.replace('\n', '').replace(' ', ''))
        total += chars
        print('  S%d  %5d자  ≈ %4.1f분' % (i, chars, chars / float(CHARS_PER_MIN)))
    prs.save(path)
    print('\n합계 %d자 ≈ %.1f분 (또박또박 %d자/분 기준)' % (total, total / float(CHARS_PER_MIN), CHARS_PER_MIN))
    print('저장: %s' % path)


if __name__ == '__main__':
    apply(sys.argv[1] if len(sys.argv) > 1 else DEFAULT)
