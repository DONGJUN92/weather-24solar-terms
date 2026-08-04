# -*- coding: utf-8 -*-
"""발표자료용 네이티브 PPT 도형 — 템플릿 팔레트만 쓴다.

이미지 캡처가 아니라 도형으로 그리는 이유
  · 확대해도 선명하고, 주최측 PC 어디서 열어도 같게 보인다
  · 템플릿 색(243A5E·E76F51·FBAE40·EFF5FB)과 정확히 맞는다
  · 앱 화면 문구를 '인용'으로 보여줄 수 있다 — 가짜 스크린샷이 아니다
앱이 그린 차트(궤도·기온 곡선·지도)는 그 자체가 증거이므로 실물 PNG를 쓴다.
"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x24, 0x3A, 0x5E)
TERRA = RGBColor(0xE7, 0x6F, 0x51)
AMBER = RGBColor(0xFB, 0xAE, 0x40)
FILL = RGBColor(0xEF, 0xF5, 0xFB)
CREAM = RGBColor(0xF3, 0xEC, 0xDC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x7A, 0x8C)
IN = 914400.0


def _tf(shape, lines, size, color, bold=False, align=PP_ALIGN.LEFT, space=0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(int(0.06 * IN))
    tf.margin_top = tf.margin_bottom = Emu(int(0.03 * IN))
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_before = Pt(space if i else 0)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return shape


def textbox(slide, x, y, w, h, lines, size=11, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, vert=None):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    sh.text_frame.vertical_anchor = anchor
    if vert:                                  # 'vert270' = 아래에서 위로 읽는 축 라벨
        sh.text_frame._bodyPr.set('vert', vert)
    return _tf(sh, lines, size, color, bold, align, space)


def box(slide, x, y, w, h, fill=WHITE, line=None, lw=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    sh.text_frame.text = ''
    return sh


def caption(slide, x, y, w, text, size=9.5):
    """도형·이미지 아래 한 줄 설명 — 무엇을 보고 있는지 밝힌다."""
    return textbox(slide, x, y, w, 0.24, [text], size=size, color=GREY)


# ─────────────────────────────────────────────────── 인용 카드
def quote_card(slide, x, y, w, h, label, lines, accent=TERRA,
               label_size=10.5, body_size=10.5):
    """앱 화면에 실제로 있는 문구를 인용하는 카드. 좌측에 강조 띠."""
    box(slide, x, y, w, h, fill=WHITE, line=RGBColor(0xC8, 0xD4, 0xE4))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(0.055), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(slide, x + 0.14, y + 0.06, w - 0.24, 0.22, [label],
            size=label_size, color=accent, bold=True)
    textbox(slide, x + 0.14, y + 0.30, w - 0.24, h - 0.36, lines,
            size=body_size, color=NAVY, space=2)


# ─────────────────────────────────────────────────── 단계 스트립
def step_strip(slide, x, y, w, h, steps, accent=NAVY, size=10, num_size=9):
    """① → ② → ③ 가로 흐름. steps = [(번호, 제목, 보조설명), ...]"""
    n = len(steps)
    gap = 0.20
    bw = (w - gap * (n - 1)) / n
    for i, (num, title, sub) in enumerate(steps):
        bx = x + i * (bw + gap)
        filled = (i == 0)
        box(slide, bx, y, bw, h, fill=(accent if filled else FILL),
            line=(None if filled else RGBColor(0xC8, 0xD4, 0xE4)))
        textbox(slide, bx + 0.06, y + 0.05, bw - 0.12, 0.20, [num],
                size=num_size, color=(CREAM if filled else accent), bold=True)
        textbox(slide, bx + 0.06, y + 0.24, bw - 0.12, 0.26, [title],
                size=size, color=(WHITE if filled else NAVY), bold=True)
        if sub:
            textbox(slide, bx + 0.06, y + 0.50, bw - 0.12, h - 0.54, [sub],
                    size=size - 1.5, color=(CREAM if filled else GREY))
        if i < n - 1:
            a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(bx + bw + 0.035), Inches(y + h / 2 - 0.055),
                                       Inches(gap - 0.07), Inches(0.11))
            a.fill.solid(); a.fill.fore_color.rgb = AMBER
            a.line.fill.background(); a.shadow.inherit = False


# ─────────────────────────────────────────────────── 계단
def ladder(slide, x, y, w, h, rungs, accent=NAVY, size=10):
    """난이도 계단 — rungs = [(라벨, 예시), ...] 왼쪽 낮고 오른쪽 높다."""
    n = len(rungs)
    gap = 0.14
    bw = (w - gap * (n - 1)) / n
    for i, (label, ex) in enumerate(rungs):
        frac = (i + 1) / float(n)
        bh = h * (0.45 + 0.55 * frac)
        bx = x + i * (bw + gap)
        by = y + (h - bh)
        shade = RGBColor(0xEF, 0xF5, 0xFB) if i == 0 else (
            RGBColor(0xCF, 0xDF, 0xEF) if i == 1 else NAVY)
        box(slide, bx, by, bw, bh, fill=shade, line=RGBColor(0xC8, 0xD4, 0xE4))
        col = WHITE if i == n - 1 else NAVY
        textbox(slide, bx + 0.06, by + 0.07, bw - 0.12, 0.24, [label],
                size=size, color=col, bold=True)
        textbox(slide, bx + 0.06, by + 0.31, bw - 0.12, bh - 0.36,
                [ex] if isinstance(ex, str) else list(ex),
                size=size - 1, color=(CREAM if i == n - 1 else GREY), space=3)


# ─────────────────────────────────────────────────── 2×2 포지셔닝
def matrix2x2(slide, x, y, w, h, xlab, ylab, cells, size=9.5):
    """cells = {(col,row): (제목, 보조)} · col 0=좌 1=우 · row 0=상 1=하"""
    ox, oy = 0.62, 0.30          # 축 라벨 자리
    gw, gh = w - ox, h - oy
    cw, ch = gw / 2.0, gh / 2.0
    for (c, r), (title, sub) in cells.items():
        bx = x + ox + c * cw
        by = y + r * ch
        hero = (c == 1 and r == 0)
        box(slide, bx + 0.03, by + 0.03, cw - 0.06, ch - 0.06,
            fill=(NAVY if hero else WHITE),
            line=(None if hero else RGBColor(0xC8, 0xD4, 0xE4)))
        textbox(slide, bx + 0.10, by + 0.09, cw - 0.20, 0.24, [title],
                size=size + 0.5, color=(WHITE if hero else NAVY), bold=True)
        textbox(slide, bx + 0.10, by + 0.33, cw - 0.20, ch - 0.40, [sub],
                size=size - 0.5, color=(CREAM if hero else GREY))
    # 축 — 세로 라벨은 vert270 로 돌려 캡션과 겹치지 않게 한다
    textbox(slide, x - 0.02, y, ox - 0.06, gh, [ylab], size=size - 0.5,
            color=TERRA, bold=True, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, vert='vert270')
    textbox(slide, x + ox, y + gh + 0.02, gw, oy - 0.04, [xlab], size=size - 0.5,
            color=TERRA, bold=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────── 네이티브 표
def mini_table(slide, x, y, w, rows, colw, row_h=0.26, head_h=0.28, size=10):
    """헤더 1행 + 본문. rows[0]=헤더. colw=열 비율 리스트."""
    ncol = len(rows[0])
    tot = float(sum(colw))
    tbl = slide.shapes.add_table(len(rows), ncol, Inches(x), Inches(y),
                                 Inches(w), Inches(head_h + row_h * (len(rows) - 1))).table
    tbl.first_row = True
    tbl.horz_banding = False
    for j, cwr in enumerate(colw):
        tbl.columns[j].width = Emu(int(w * IN * cwr / tot))
    tbl.rows[0].height = Emu(int(head_h * IN))
    for i in range(1, len(rows)):
        tbl.rows[i].height = Emu(int(row_h * IN))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = cell.margin_right = Emu(int(0.05 * IN))
            cell.margin_top = cell.margin_bottom = Emu(int(0.01 * IN))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if i == 0 else (
                RGBColor(0xFF, 0xF3, 0xE2) if val.startswith('!') else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val[1:] if val.startswith('!') else val
            r.font.size = Pt(size)
            r.font.bold = (i == 0) or val.startswith('!')
            r.font.color.rgb = CREAM if i == 0 else (TERRA if val.startswith('!') else NAVY)
    return tbl


# ─────────────────────────────────────────────────── 지표 배지
def stat_row(slide, x, y, w, h, stats, size=13, lab=8.5):
    """큰 숫자 + 라벨 묶음. stats = [(값, 라벨), ...]"""
    n = len(stats)
    gap = 0.12
    bw = (w - gap * (n - 1)) / n
    for i, (val, label) in enumerate(stats):
        bx = x + i * (bw + gap)
        box(slide, bx, y, bw, h, fill=FILL, line=RGBColor(0xC8, 0xD4, 0xE4))
        textbox(slide, bx + 0.04, y + 0.05, bw - 0.08, h * 0.52, [val],
                size=size, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, bx + 0.04, y + h * 0.55, bw - 0.08, h * 0.42, [label],
                size=lab, color=GREY, align=PP_ALIGN.CENTER)
